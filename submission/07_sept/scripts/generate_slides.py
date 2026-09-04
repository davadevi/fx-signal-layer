"""
Generate FX Signal Layer hackathon presentation.
Run: python generate_slides.py
Output: slides.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colors ──────────────────────────────────────────────────────────────────
RED    = RGBColor(0xCC, 0x00, 0x00)   # Alfa-Bank red
DARK   = RGBColor(0x1A, 0x1A, 0x1A)  # Near-black
GRAY   = RGBColor(0x55, 0x55, 0x55)  # Body gray
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)  # Table alt row
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED_BG = RGBColor(0xF8, 0xE8, 0xE8)  # Light red for header rows
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
AMBER  = RGBColor(0xE6, 0x5C, 0x00)

W = 13.33  # slide width inches (16:9)
H = 7.5   # slide height inches


# ── Helpers ─────────────────────────────────────────────────────────────────

def make_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(layout)
    # white background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


def add_box(slide, x, y, w, h, text, size=18, color=DARK, bold=False,
            italic=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text        = text
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.bold   = bold
    run.font.italic = italic
    run.font.name   = "Calibri"
    p.alignment     = align
    return txb


def add_multiline(slide, x, y, w, h, lines, size=14, color=DARK,
                  bold=False, line_spacing=None):
    """lines: list of str or (str, dict) for per-line overrides."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if isinstance(line, tuple):
            text, opts = line
        else:
            text, opts = line, {}
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(opts.get("size", size))
        run.font.color.rgb = opts.get("color", color)
        run.font.bold  = opts.get("bold", bold)
        run.font.name  = "Calibri"
        if line_spacing:
            from pptx.util import Pt as UPt
            p.line_spacing = UPt(line_spacing)
    return txb


def add_rule(slide, x, y, w, color=RED, thickness=3):
    """Horizontal divider line."""
    from pptx.util import Emu
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.LINE = 1, but we use connector
        Inches(x), Inches(y), Inches(w), Emu(thickness * 9144)
    )
    line.line.color.rgb = color
    line.line.width = Emu(thickness * 9144)
    return line


def slide_header(slide, title: str, subtitle: str = ""):
    """Standard slide header: red title + gray subtitle + red underline."""
    add_box(slide, 0.5, 0.25, W - 1.0, 0.65, title,
            size=28, color=RED, bold=True)
    if subtitle:
        add_box(slide, 0.5, 0.85, W - 1.0, 0.4, subtitle,
                size=13, color=GRAY)
    # red rule
    from pptx.util import Emu
    line = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(0.5), Inches(0.92 if not subtitle else 1.22),
        Inches(W - 0.5), Inches(0.92 if not subtitle else 1.22),
    )
    line.line.color.rgb = RED
    line.line.width = Emu(18000)


def add_table(slide, headers, rows, x, y, w, h,
              col_widths=None, header_color=RED_BG):
    """Add a table. headers: list[str], rows: list[list[str]]."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(x), Inches(y), Inches(w), Inches(h)
    ).table

    # col widths
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    def set_cell(cell, text, bold=False, color=DARK, bg=None,
                 align=PP_ALIGN.LEFT, size=12):
        tf = cell.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.name  = "Calibri"
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # header row
    for j, h_text in enumerate(headers):
        set_cell(tbl.cell(0, j), h_text, bold=True, bg=header_color,
                 align=PP_ALIGN.CENTER, size=12)

    # data rows
    for i, row in enumerate(rows):
        bg = LGRAY if i % 2 == 1 else WHITE
        for j, cell_text in enumerate(row):
            set_cell(tbl.cell(i + 1, j), str(cell_text),
                     bg=bg, align=PP_ALIGN.LEFT, size=11)

    return tbl


# ── Slides ───────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)

    # Red accent bar left
    bar = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0), Inches(0), Inches(0.4), Inches(H)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

    # Main title
    add_box(slide, 0.8, 1.8, 11.5, 1.2,
            "FX Signal Layer",
            size=48, color=DARK, bold=True)

    # Subtitle
    add_box(slide, 0.8, 3.0, 11.5, 0.7,
            "Система push-сигналов для трансграничных переводов",
            size=22, color=GRAY)

    # Divider
    from pptx.util import Emu
    line = slide.shapes.add_connector(
        1, Inches(0.8), Inches(3.75), Inches(11.5), Inches(3.75)
    )
    line.line.color.rgb = RED
    line.line.width = Emu(27000)

    # Details
    add_box(slide, 0.8, 4.0, 11.5, 0.5,
            "Alfa-Bank Hackathon · Сентябрь 2026",
            size=16, color=GRAY)

    add_box(slide, 0.8, 4.7, 11.5, 0.5,
            "Давид Гусейнов · София · Варя",
            size=15, color=DARK, bold=True)


def slide_problem(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Проблема")

    points = [
        "Клиент переводит деньги когда удобно ему, а не когда выгоден курс — разница может достигать 2–3% на горизонте недели.",
        "В приложении нет сигнала о выгодности текущего момента: пользователь не имеет ориентира кроме собственной интуиции.",
        "Трудовые мигранты совершают 1–2 перевода в месяц, и ошибка в выборе дня стоит реальных денег — при переводе 30 000 ₽ это 600–900 ₽.",
    ]
    icons = ["01", "02", "03"]
    y_positions = [1.45, 2.95, 4.45]

    for i, (icon, point, y) in enumerate(zip(icons, points, y_positions)):
        # number circle box
        num = slide.shapes.add_shape(
            1, Inches(0.5), Inches(y), Inches(0.55), Inches(0.55)
        )
        num.fill.solid()
        num.fill.fore_color.rgb = RED
        num.line.fill.background()
        tf = num.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = icon
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"

        add_box(slide, 1.25, y - 0.05, 11.5, 0.75, point,
                size=15, color=DARK)


def slide_task(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Постановка задачи")

    add_box(slide, 0.5, 1.35, 12.3, 0.6,
            "Разработать систему, которая определяет выгодные моменты для трансграничных переводов "
            "и генерирует push-уведомление клиенту с описанием текущей ситуации на основе фактов прошлого.",
            size=14, color=DARK)

    headers = ["Метрика", "Целевое значение", "Обоснование"]
    rows = [
        ["Lift над случайным днём", "≥ 1.3", "Порог экономической значимости сигнала"],
        ["Частота сигналов", "1–2 / коридор / неделю", "Достаточно для покрытия клиентских паттернов"],
        ["Горизонт оценки", "h = 5 рабочих дней", "Реалистичный лаг принятия решения клиентом"],
        ["Покрытие коридоров", "5 (TJS, UZS, KGS, AMD, KZT)", "Ключевые направления трудовых переводов"],
        ["Ошибка асимметрии", "FP > FN (вес FP = 3×)", "Ложный сигнал дороже пропущенного"],
    ]
    add_table(slide, headers, rows,
              x=0.5, y=2.05, w=12.3, h=4.6,
              col_widths=[3.5, 3.3, 5.5])


def slide_product(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Продуктовый фокус",
                 "Пользователь · Отличия · Сценарий")

    col_w = 4.0
    col_gap = 0.15
    col_y = 1.5
    col_h = 5.7
    x1, x2, x3 = 0.5, 4.65, 8.8

    # ── Column backgrounds ────────────────────────────────────────────────────
    for x in (x1, x2, x3):
        bg = slide.shapes.add_shape(
            1, Inches(x), Inches(col_y), Inches(col_w), Inches(col_h)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = LGRAY
        bg.line.color.rgb = RED

    # ── Block 1: Кто и зачем ─────────────────────────────────────────────────
    add_box(slide, x1 + 0.15, col_y + 0.15, col_w - 0.3, 0.4,
            "Блок 1 — Кто и зачем", size=12, color=RED, bold=True)

    add_box(slide, x1 + 0.15, col_y + 0.65, col_w - 0.3, 0.65,
            "Трудовой мигрант · зарплатный цикл · один коридор / один получатель\n"
            "Часовой пояс UTC+5 → запрет пуша 22:00–08:00 по времени получателя",
            size=10, color=DARK, bold=False)

    # JTBD quote box
    quote_box = slide.shapes.add_shape(
        1, Inches(x1 + 0.15), Inches(col_y + 1.3), Inches(col_w - 0.3), Inches(3.4)
    )
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = WHITE
    quote_box.line.color.rgb = RED

    tf = quote_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "JTBD"
    r0.font.size = Pt(10)
    r0.font.bold = True
    r0.font.color.rgb = RED
    r0.font.name = "Calibri"

    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    r1.text = ("« Когда мне нужно отправить деньги семье, "
               "я хочу понимать — сегодня хороший день по курсу "
               "или лучше подождать, — чтобы до получателя дошло "
               "больше денег без необходимости ежедневно следить "
               "за валютным рынком самому. »")
    r1.font.size = Pt(12)
    r1.font.italic = True
    r1.font.color.rgb = DARK
    r1.font.name = "Calibri"

    # ── Block 2: Чем отличаемся ───────────────────────────────────────────────
    add_box(slide, x2 + 0.15, col_y + 0.15, col_w - 0.3, 0.4,
            "Блок 2 — Чем отличаемся", size=12, color=RED, bold=True)

    competitors = [
        (
            "Курс-трекеры (XE, Wise alerts)",
            "Алерт на порог, заданный пользователем вручную. "
            "У нас: момент определяет сигнальный слой — клиенту не нужно настраивать пороги.",
        ),
        (
            "«Best time to send» у Wise / Remitly",
            "Исторический контекст без walk-forward валидации. "
            "У нас: ML + walk-forward OOT + явный запрет на прогноз.",
        ),
        (
            "E-commerce «успей купить»",
            "Срочность = имплицитный прогноз. "
            "У нас: только факты о прошлом/настоящем, никогда не предсказание.",
        ),
    ]

    comp_y = col_y + 0.68
    for (comp_title, comp_body) in competitors:
        # small red marker
        marker = slide.shapes.add_shape(
            1, Inches(x2 + 0.15), Inches(comp_y), Inches(0.08), Inches(0.85)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = RED
        marker.line.fill.background()

        add_box(slide, x2 + 0.33, comp_y, col_w - 0.5, 0.35,
                comp_title, size=11, color=DARK, bold=True)
        add_box(slide, x2 + 0.33, comp_y + 0.35, col_w - 0.5, 0.6,
                comp_body, size=10, color=GRAY)
        comp_y += 1.6

    # ── Block 3: Сценарий ─────────────────────────────────────────────────────
    add_box(slide, x3 + 0.15, col_y + 0.15, col_w - 0.3, 0.4,
            "Блок 3 — Сценарий", size=12, color=RED, bold=True)

    steps = [
        ("PUSH", "Пуш с фактом о курсе", RED),
        ("ЭКРАН", "Экран текущего курса", RGBColor(0x15, 0x65, 0xC0)),
        ("ФОРМА", "Предзаполненная форма перевода", GREEN),
    ]

    step_y = col_y + 0.72
    for i, (icon_text, step_label, icon_color) in enumerate(steps):
        # icon circle
        icon_box = slide.shapes.add_shape(
            1, Inches(x3 + 0.3), Inches(step_y), Inches(0.85), Inches(0.72)
        )
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = icon_color
        icon_box.line.fill.background()

        tf_i = icon_box.text_frame
        pi = tf_i.paragraphs[0]
        pi.alignment = PP_ALIGN.CENTER
        ri = pi.add_run()
        ri.text = icon_text
        ri.font.size = Pt(9)
        ri.font.bold = True
        ri.font.color.rgb = WHITE
        ri.font.name = "Calibri"

        add_box(slide, x3 + 1.3, step_y + 0.1, col_w - 1.5, 0.6,
                step_label, size=12, color=DARK)

        # down arrow (not after last)
        if i < len(steps) - 1:
            from pptx.util import Emu
            arr = slide.shapes.add_connector(
                1,
                Inches(x3 + 0.72), Inches(step_y + 0.72),
                Inches(x3 + 0.72), Inches(step_y + 1.2),
            )
            arr.line.color.rgb = RED
            arr.line.width = Emu(18000)

        step_y += 1.52


def slide_approach(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Выбранный подход",
                 "Log-return percentile rank как основной индикатор")

    # Pipeline flow boxes
    steps = [
        ("CBR XML", "Источник данных\nЦБ РФ"),
        ("Нормализация", "VunitRate,\nforward-fill"),
        ("Индикаторы", "Log-return\npercentile rank"),
        ("Бэктест", "Walk-forward,\nCI bootstrap"),
        ("ML-слой", "LightGBM,\nasymm. loss"),
        ("Pipeline", "Cooldown,\ntier filter"),
        ("Push-текст", "Compliance\nvalidated"),
    ]

    box_w = 1.55
    box_h = 1.1
    start_x = 0.35
    y_top = 1.5
    gap = 0.15

    for i, (title, desc) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        box = slide.shapes.add_shape(
            1, Inches(x), Inches(y_top), Inches(box_w), Inches(box_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LGRAY
        box.line.color.rgb = RED

        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = title
        r0.font.size = Pt(11)
        r0.font.bold = True
        r0.font.color.rgb = RED
        r0.font.name = "Calibri"

        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = desc
        r1.font.size = Pt(9)
        r1.font.color.rgb = GRAY
        r1.font.name = "Calibri"

        # arrow (except after last)
        if i < len(steps) - 1:
            ax = x + box_w + 0.01
            arrow = slide.shapes.add_connector(
                1,
                Inches(ax), Inches(y_top + box_h / 2),
                Inches(ax + gap), Inches(y_top + box_h / 2),
            )
            arrow.line.color.rgb = RED
            from pptx.util import Emu
            arrow.line.width = Emu(18000)

    # Key decision callout
    add_box(slide, 0.5, 2.85, 12.3, 0.55,
            "Ключевое решение: индикатор работает на log-return (I(0) — стационарный ряд), "
            "а не на абсолютном уровне курса (I(1) — требует дифференцирования).",
            size=13, color=DARK, italic=True)

    add_box(slide, 0.5, 3.5, 12.3, 1.5,
            "Почему это важно: абсолютный уровень курса имеет единичный корень — percentile rank по нему "
            "технически некорректен и даёт lift < 1.0. Log-return является стационарным, "
            "percentile rank интерпретируется корректно, lift на KGS/TJS достигает 1.6.",
            size=13, color=GRAY)

    # Compliance note
    box = slide.shapes.add_shape(
        1, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.85)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RED_BG
    box.line.color.rgb = RED

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Compliance ограничение: push-текст содержит только факты прошлого и настоящего. "
    r.font.size = Pt(12)
    r.font.color.rgb = DARK
    r.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = "Запрещены: предсказания, срочность, гарантии, инвестиционные формулировки."
    r2.font.size = Pt(12)
    r2.font.bold = True
    r2.font.color.rgb = RED
    r2.font.name = "Calibri"


def slide_architecture(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Архитектура решения",
                 "Модульная структура · 2 150 строк production-кода")

    modules = [
        ("src/data/",       "Загрузка и нормализация курсов\nCBR XML → parquet, VunitRate, forward-fill",                               0.4, 1.45),
        ("src/indicators/", "8 индикаторов · BaseIndicator\nLog-return percentile, RSI, Volatility regime, Bollinger, Calendar",        4.4, 1.45),
        ("src/backtest/",   "Walk-forward engine\nHit rate A/B, lift, CI bootstrap (2000 resamples, 90d blocks)",                       8.6, 1.45),
        ("src/ml/",         "LightGBM walk-forward\nFeatures (10 cols), labels (local min h=5), asymmetric loss FP×3, two-tier thresh", 0.4, 3.65),
        ("src/pipeline/",   "Генерация сигналов\nMandatory / optional tier, cooldown 3d, cap 2 sig/wk, compliance check",              4.4, 3.65),
        ("src/texts/",      "Push-текст\nCompliance validator + formatter, шаблоны только по фактам",                                   8.6, 3.65),
    ]

    box_w, box_h = 3.8, 1.7

    for (name, desc, x, y) in modules:
        box = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(box_w), Inches(box_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LGRAY
        box.line.color.rgb = RED

        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.LEFT
        r0 = p0.add_run()
        r0.text = name
        r0.font.size = Pt(13)
        r0.font.bold = True
        r0.font.color.rgb = RED
        r0.font.name = "Calibri"

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = desc
        r1.font.size = Pt(10)
        r1.font.color.rgb = GRAY
        r1.font.name = "Calibri"

    # Data flow arrows (top row: data→indicators→backtest)
    from pptx.util import Emu
    arrow_y = Inches(2.3)
    for ax in [Inches(4.2), Inches(8.2)]:
        conn = slide.shapes.add_connector(
            1, ax, arrow_y, ax + Inches(0.2), arrow_y
        )
        conn.line.color.rgb = RED
        conn.line.width = Emu(18000)

    # Bottom row arrows: ml→pipeline→texts
    arrow_y2 = Inches(4.5)
    for ax in [Inches(4.2), Inches(8.2)]:
        conn = slide.shapes.add_connector(
            1, ax, arrow_y2, ax + Inches(0.2), arrow_y2
        )
        conn.line.color.rgb = RED
        conn.line.width = Emu(18000)

    # Stats strip
    add_box(slide, 0.4, 5.55, 12.5, 0.45,
            "Данные: 19 500 строк · 5 коридоров · 1 642 торговых дня/коридор · 6 ADR · 3 итерационных лога · 25 unit-тестов",
            size=11, color=GRAY)


def slide_stages(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Выполненные этапы")

    headers = ["Этап", "Статус", "Ключевое решение", "Результат"]
    rows = [
        ["01 · Данные",      "✓ Завершён",   "CBR VunitRate; cutoff 2022-04-01 (post-shock); forward-fill выходных",
         "19 500 строк, parquet, is_trading_day"],
        ["02 · Индикаторы",  "✓ Завершён",   "Log-return rank (I(0)) вместо абсолютного уровня (I(1))",
         "8 вариантов; LogReturnPercentile — основной"],
        ["03 · Бэктест",     "✓ Завершён",   "Walk-forward 2y/3m, embargo 5d, hit def A+B, CI bootstrap",
         "KGS/TJS: lift 2.1–2.5, CI lower bound > 1.0 на OOT"],
        ["04 · ML-слой",     "✓ Завершён",   "LightGBM, асимметричный loss (FP × 3); деградирует vs. чистый индикатор",
         "Из pipeline исключён — чистый LogReturnPercentile стабильнее"],
        ["05 · Pipeline",    "✓ Завершён",   "confirm_days=2, stateful cooldown, compliance validator",
         "generate_signals() + stateful history; 30 unit-тестов"],
        ["06 · OOT-валидация", "✓ Завершён", "Pathwise CI bootstrap; OOT split 2025-07-01",
         "3 коридора CI > 1.0: KGS h=5,10; TJS h=5,10; AMD h=10"],
    ]
    add_table(slide, headers, rows,
              x=0.5, y=1.45, w=12.3, h=5.8,
              col_widths=[2.0, 1.5, 4.8, 4.0])


def slide_results(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Результаты — OOT-валидация",
                 "LogReturnPercentile · confirm_days=2 · walk-forward 2y/3m · OOT с 2025-07-01")

    headers = ["Коридор", "h", "IS Lift", "OOT Lift", "CI 95% ↓", "Сигн./нед", "Статус"]
    rows = [
        ["RUB / KGS (сом)",    "5",    "2.10", "1.84", "1.36", "0.057", "✓ CI > 1.0"],
        ["RUB / KGS (сом)",    "10",   "2.17", "2.24", "1.15", "0.057", "✓ CI > 1.0"],
        ["RUB / TJS (сомони)", "5",    "2.35", "1.76", "1.76", "0.065", "✓ CI > 1.0"],
        ["RUB / TJS (сомони)", "10",   "2.52", "2.13", "1.65", "0.065", "✓ CI > 1.0"],
        ["RUB / AMD (драм)",   "10",   "2.81", "2.16", "1.30", "0.033", "✓ CI > 1.0"],
        ["RUB / UZS (сум)",    "5–20", "1.2–1.9", "1.0–1.3", "< 1.0", "0.065", "✗ CI не проходит"],
        ["RUB / KZT (тенге)",  "все",  "~1.1", "NaN",  "< 1.0", "0.033", "✗ Нет OOT сигналов"],
    ]
    add_table(slide, headers, rows,
              x=0.5, y=1.45, w=9.0, h=4.8,
              col_widths=[2.2, 0.5, 1.0, 1.1, 1.0, 1.1, 2.1])

    # Callout: key findings
    box = slide.shapes.add_shape(
        1, Inches(9.8), Inches(1.45), Inches(3.0), Inches(4.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RED_BG
    box.line.color.rgb = RED

    tf = box.text_frame
    tf.word_wrap = True

    lines = [
        ("3 коридора с CI > 1.0", {"size": 13, "bold": True, "color": RED}),
        ("", {}),
        ("KGS и TJS: OOT lift выше IS на h=10 — сигнал не деградирует на новых данных.", {"size": 10, "color": DARK}),
        ("", {}),
        ("AMD: проходит только на h=10 (N=4 — малая выборка).", {"size": 10, "color": AMBER}),
        ("", {}),
        ("KZT: нет сигналов в OOT периоде — индикатор не работает.", {"size": 10, "color": GRAY}),
        ("", {}),
        ("Частота: 0.033–0.065 сигн./нед — ниже цели 1–2.", {"size": 10, "color": DARK}),
    ]
    first = True
    for (text, opts) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size  = Pt(opts.get("size", 10))
        r.font.bold  = opts.get("bold", False)
        r.font.color.rgb = opts.get("color", DARK)
        r.font.name  = "Calibri"

    add_box(slide, 0.5, 6.45, 12.3, 0.45,
            "Базовая модель (абсолютный percentile) — lift 0.86–0.97 (< 1.0): работает хуже случайного выбора дня.",
            size=11, color=GRAY, italic=True)


def slide_open_questions(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Открытые вопросы")

    questions = [
        (
            "Частота сигналов",
            "confirm_days=2 даёт 0.05–0.07 сигн./нед — в 15–30× ниже цели 1–2/нед. "
            "Relaxed вариант (confirm_days=0) даёт ~0.46 сигн./нед, но CI не проходит 1.0 на OOT. "
            "Продуктовый выбор: редкий и надёжный vs. частый и менее точный."
        ),
        (
            "bps отрицательный на h=5",
            "Hit rate > 50% подтверждён на KGS/TJS, но средняя выгода в базисных пунктах отрицательна на h=5 "
            "— сигнал угадывает направление, но приходит после выгодного движения. На h=10/20 bps улучшается."
        ),
        (
            "AMD — малая выборка",
            "AMD проходит CI на h=10 (lift 2.81, OOT 2.16, CI low 1.30), но N=4 сигнала. "
            "Статистическая устойчивость требует накопления OOT-сигналов в реальном пилоте."
        ),
    ]

    for i, (q_title, q_body) in enumerate(questions):
        y = 1.45 + i * 1.8
        # title
        add_box(slide, 0.5, y, 12.3, 0.45, f"{i+1}. {q_title}",
                size=15, color=RED, bold=True)
        add_box(slide, 0.8, y + 0.42, 12.0, 0.95, q_body,
                size=13, color=DARK)


def slide_risks(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Ключевые риски")

    headers = ["Риск", "Вероятность", "Митигация"]
    rows = [
        [
            "Частота < 1 сигн./нед после полной OOT валидации на нетронутых данных",
            "Высокая",
            "Ввести weak-tier (confirm_days=0) с пониженным весом и separate cooldown; "
            "или добавить calendar seasonality как второй источник сигналов"
        ],
        [
            "Режимный сдвиг (санкции 2022, внешние шоки) нарушает стабильность обучения модели",
            "Средняя",
            "Cutoff 2022-04-01 устраняет шок из обучения; walk-forward + embargo 5d "
            "предотвращает lookahead; volatility regime filter глушит сигналы в кризис"
        ],
        [
            "Push-текст нарушает compliance (неявное предсказание, срочность)",
            "Низкая",
            "Встроенный compliance_validator проверяет каждый текст; все шаблоны содержат "
            "только факты прошлого ('укрепился', 'выгоднее чем в X% дней')"
        ],
        [
            "Переобучение ML-слоя на in-sample данных без истинного OOT",
            "Средняя",
            "OOT split зафиксирован на 2025-07-01; тест окна не пересекаются с "
            "отбором параметров; планируется финальная OOT прогонка"
        ],
    ]
    add_table(slide, headers, rows,
              x=0.5, y=1.45, w=12.3, h=5.5,
              col_widths=[4.5, 1.8, 6.0])


def slide_plan(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Что сделано · 07 сентября")

    done_items = [
        ("✓ OOT-валидация с pathwise CI bootstrap",
         "3 коридора проходят CI > 1.0; KGS/TJS стабильны на нетронутых данных"),
        ("✓ Синхронизация backtest ↔ runtime",
         "generate_signals() использует ту же логику, что walk-forward engine"),
        ("✓ Stateful cooldown",
         "История сигналов сохраняется в data/signal_history.json; cooldown работает между запусками"),
        ("✓ REPRODUCE.md · LIMITATIONS.md",
         "Команды воспроизведения; честные ограничения: частота, bps, малая выборка AMD"),
        ("✓ Pilot Design (05_pilot_design.md)",
         "A/B дизайн, первичные метрики, стоп-критерии, оценка длительности 3–6 мес."),
    ]

    y = 1.45
    for title, subtitle in done_items:
        # green left border strip
        strip = slide.shapes.add_shape(
            1, Inches(0.5), Inches(y), Inches(0.12), Inches(0.82)
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = GREEN
        strip.line.fill.background()

        add_box(slide, 0.75, y, 11.8, 0.42, title,
                size=13, color=DARK, bold=True)
        add_box(slide, 0.75, y + 0.40, 11.8, 0.38, subtitle,
                size=11, color=GRAY)
        y += 1.0

    # Final line
    add_box(slide, 0.5, 6.55, 12.3, 0.6,
            "→ 07 сентября: Финальная презентация + Pilot Design",
            size=14, color=RED, bold=True)


def slide_team(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Команда")

    members = [
        {
            "name":   "Давид Гусейнов",
            "role":   "ML Engineer",
            "tasks":  [
                "Данные: загрузка CBR, нормализация, parquet",
                "Индикаторы: 8 вариантов, LogReturnPercentile",
                "Бэктест: walk-forward engine, CI bootstrap",
                "ML: LightGBM, asymmetric loss, features/labels",
            ],
            "pct":    "33%",
        },
        {
            "name":   "София",
            "role":   "Product Manager",
            "tasks":  [
                "Продуктовая гипотеза и портрет пользователя",
                "ADR: 6 файлов архитектурных решений",
                "Compliance-требования к push-текстам",
                "Документация: интерфейсы, PIPELINE.md",
            ],
            "pct":    "33%",
        },
        {
            "name":   "Варя",
            "role":   "Data Engineer",
            "tasks":  [
                "Pipeline: generate_signals(), mandatory/optional",
                "Push-текст: шаблоны + compliance validator",
                "Тесты: 25 unit-тестов, no-lookahead suite",
                "Валидация: lift tables, text report",
            ],
            "pct":    "33%",
        },
    ]

    col_x = [0.5, 4.6, 8.7]
    box_w = 3.9
    box_h = 5.2

    for i, (member, x) in enumerate(zip(members, col_x)):
        # Card background
        card = slide.shapes.add_shape(
            1, Inches(x), Inches(1.45), Inches(box_w), Inches(box_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LGRAY
        card.line.color.rgb = RED

        # Name header
        name_box = slide.shapes.add_shape(
            1, Inches(x), Inches(1.45), Inches(box_w), Inches(0.75)
        )
        name_box.fill.solid()
        name_box.fill.fore_color.rgb = RED
        name_box.line.fill.background()

        tf = name_box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = member["name"]
        r0.font.size = Pt(14)
        r0.font.bold = True
        r0.font.color.rgb = WHITE
        r0.font.name = "Calibri"

        # Role
        add_box(slide, x + 0.15, 2.3, box_w - 0.3, 0.4,
                member["role"], size=13, color=RED, bold=True)

        # Tasks
        task_text = "\n".join(f"· {t}" for t in member["tasks"])
        add_box(slide, x + 0.15, 2.75, box_w - 0.3, 3.0,
                task_text, size=11, color=DARK)

        # Participation badge
        pct_box = slide.shapes.add_shape(
            1, Inches(x + box_w - 0.9), Inches(6.15), Inches(0.75), Inches(0.35)
        )
        pct_box.fill.solid()
        pct_box.fill.fore_color.rgb = RED
        pct_box.line.fill.background()

        tf2 = pct_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = member["pct"]
        r2.font.size = Pt(11)
        r2.font.bold = True
        r2.font.color.rgb = WHITE
        r2.font.name = "Calibri"


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = make_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_task(prs)
    slide_product(prs)
    slide_approach(prs)
    slide_architecture(prs)
    slide_stages(prs)
    slide_results(prs)
    slide_open_questions(prs)
    slide_risks(prs)
    slide_plan(prs)
    slide_team(prs)

    out = os.path.join(os.path.dirname(__file__), "..", "01_presentation.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
